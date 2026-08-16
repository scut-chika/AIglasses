# -*- coding: utf-8 -*-
"""生成《暖眸 · 居家养老 AI 眼镜智能体》项目说明 PPT（文档向，约 20 页）。

运行：py make_ppt.py
输出：暖眸-居家养老AI眼镜智能体-项目说明.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 主题 ----------------
BG = RGBColor(0x0D, 0x10, 0x20)
PANEL = RGBColor(0x15, 0x1A, 0x30)
PANEL2 = RGBColor(0x1A, 0x20, 0x40)
BORDER = RGBColor(0x2A, 0x31, 0x57)
TEXT = RGBColor(0xE8, 0xEA, 0xF6)
MUTED = RGBColor(0x8B, 0x93, 0xB8)
ACCENT = RGBColor(0x6D, 0x8D, 0xFF)
AMBER = RGBColor(0xFF, 0xB4, 0x54)
GREEN = RGBColor(0x57, 0xD9, 0xA3)
RED = RGBColor(0xFF, 0x6B, 0x81)
PER = RGBColor(0x5A, 0xB0, 0xFF)
LOG = RGBColor(0xCD, 0xD3, 0xEC)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_font(run, size=14, bold=False, color=TEXT, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        set_font(r, size, bold, color)
    return box


def add_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=0.12, line_w=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def shape_text(shape, text, size=13, bold=False, color=TEXT,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = line
        set_font(r, size, bold, color)


def add_arrow(slide, x, y, w, h, color=ACCENT, direction="right"):
    kind = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_footer(slide, page):
    add_text(slide, 10.3, 7.12, 2.6, 0.3,
             f"暖眸 · 居家养老 AI 眼镜智能体  {page:02d}",
             size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_header(slide, page, chapter, title, subtitle=None):
    add_text(slide, 0.6, 0.28, 0.9, 0.55, f"{page:02d}", size=23, bold=True, color=ACCENT)
    add_text(slide, 1.25, 0.24, 11.4, 0.35, chapter, size=11.5, bold=True, color=AMBER)
    add_text(slide, 1.25, 0.56, 11.4, 0.6, title, size=24, bold=True)
    if subtitle:
        add_text(slide, 1.27, 1.12, 11.2, 0.4, subtitle, size=12.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.1), Inches(12.13), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    line.shadow.inherit = False
    add_footer(slide, page)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def flow_slide(page, chapter, title, subtitle, steps, branches, notes):
    """闭环流程页：横向步骤 + 异常分支面板。"""
    s = new_slide()
    add_header(s, page, chapter, title, subtitle)
    n = len(steps)
    gap = 0.32
    box_w = (12.13 - gap * (n - 1)) / n
    box_h = 1.9
    y = 1.6
    x = 0.6
    for i, (t, d) in enumerate(steps):
        bx = x + i * (box_w + gap)
        box = add_rect(s, bx, y, box_w, box_h, fill=PANEL2, line=ACCENT, radius=0.1)
        shape_text(box, f"{t}\n\n{d}", size=12.5 if len(d) > 14 else 13)
        if i < n - 1:
            add_arrow(s, bx + box_w + 0.045, y + box_h / 2 - 0.12, gap - 0.09, 0.24,
                      color=ACCENT, direction="right")
    add_text(s, 0.6, 3.8, 12.1, 0.35, "异常分支 / 边界", size=15, bold=True, color=AMBER)
    for i, (t, d, c) in enumerate(branches):
        bw = (12.13 - 0.3) / len(branches)
        bbox = add_rect(s, 0.6 + i * (bw + 0.1), 4.25, bw, 1.5, fill=PANEL, line=BORDER)
        shape_text(bbox, f"{t}\n{d}", size=12.5, color=c)
    add_text(s, 0.6, 6.05, 12.1, 0.6, "闭环价值：感知 → 决策 → 行动 → 反馈 → 记忆沉淀，全链路可回溯。", size=13, color=MUTED)
    add_footer(s, page)
    set_notes(s, notes)
    return s


# ================= 01 封面 =================
s = new_slide()
add_rect(s, 0.6, 1.2, 0.14, 4.5, fill=ACCENT, line=ACCENT, radius=0.5)
add_text(s, 1.15, 1.45, 11.4, 1.0, "暖眸 · 居家养老 AI 眼镜智能体", size=44, bold=True)
add_text(s, 1.15, 2.5, 11.4, 0.5, "—— 项目说明 ——", size=20, color=AMBER)
add_text(s, 1.15, 3.35, 11.4, 0.45, "赛题方向：AI+眼镜 · 下一代个人智能终端原生智能体", size=14, color=TEXT)
add_text(s, 1.15, 3.8, 11.4, 0.45, "赛道：生活健康    重点验证：第一视角 / 语音交互 / 实时感知 / 随身陪伴 / 轻量反馈", size=13.5, color=MUTED)
add_text(s, 1.15, 4.55, 11.4, 0.4, "作品形式：真实可运行 Agent 应用（用户提供 AI Key 即可运行）", size=14, bold=True, color=GREEN)
add_text(s, 1.15, 5.0, 11.4, 0.4, "＋ Web 概念演示（无需 Key）＋ 功能闭环演示 ＋ 设计文档", size=13, color=MUTED)
add_text(s, 1.15, 6.35, 11.4, 0.4, "队伍：＿＿＿＿＿＿   |   2026 年 8 月", size=12.5, color=MUTED)
set_notes(s, "本 PPT 为项目说明文档，按「赛题理解→需求→设计→实现→验证→交付」组织。")

# ================= 02 目录 =================
s = new_slide()
add_header(s, 2, "目录", "项目说明结构")
items = [
    ("1", "赛题理解与作品定位", "五个验证点的落点、为什么是眼镜"),
    ("2", "目标用户与场景痛点", "独居老人、子女、社区护工"),
    ("3", "需求分析", "功能需求与非功能需求"),
    ("4", "功能设计", "7 项能力、P0/P1 分层、明确不做"),
    ("5", "核心任务闭环", "找物 / 用药 / 出门 / 跌倒 四个闭环"),
    ("6", "Agent 设计与架构", "主动触发引擎、分层架构、记忆系统"),
    ("7", "技术实现与产品形态", "真实版 app/、系统使用、双端 App 构想"),
    ("8", "数据与合规边界", "授权、隐私、风险、行业边界"),
    ("9", "Demo 与验证", "真实运行演示、端到端测试"),
    ("10", "迭代计划与交付物", "版本路线、参赛交付清单"),
]
for i, (num, t, d) in enumerate(items):
    r, c = divmod(i, 2)
    box = add_rect(s, 0.6 + c * 6.16, 1.45 + r * 1.28, 5.96, 1.1, fill=PANEL, line=BORDER, radius=0.1)
    add_text(s, 0.85 + c * 6.16, 1.62 + r * 1.28, 0.6, 0.5, num, size=22, bold=True, color=ACCENT)
    add_text(s, 1.55 + c * 6.16, 1.6 + r * 1.28, 4.8, 0.45, t, size=15, bold=True)
    add_text(s, 1.55 + c * 6.16, 2.02 + r * 1.28, 4.8, 0.4, d, size=11.5, color=MUTED)
add_footer(s, 2)
set_notes(s, "本页为全文目录，共 10 章。")

# ================= 03 赛题理解与作品定位 =================
s = new_slide()
add_header(s, 3, "1 赛题理解与作品定位", "AI+眼镜 · 下一代个人智能终端原生智能体")
add_text(s, 0.6, 1.5, 12.1, 0.4, "五个重点验证点在作品中的落点：", size=15, bold=True)
rows = [
    ("第一视角", "摄像头 / 照片作为感知输入，视觉模型识别“老人看到的”", ACCENT),
    ("语音交互", "浏览器语音识别输入 + 语音朗读反馈，无需按键", GREEN),
    ("实时感知", "场景 / 人物状态 / 可见物体逐轮识别，驱动决策", PER),
    ("随身陪伴", "持续记忆 + 主动触发，越用越懂老人", AMBER),
    ("轻量反馈", "一句话语音反馈，低打扰原则", RED),
]
for i, (t, d, c) in enumerate(rows):
    box = add_rect(s, 0.6, 2.0 + i * 0.62, 12.13, 0.5, fill=PANEL, line=BORDER, radius=0.12)
    add_text(s, 0.85, 2.11 + i * 0.62, 1.7, 0.35, t, size=13.5, bold=True, color=c)
    add_text(s, 2.7, 2.11 + i * 0.62, 9.8, 0.35, d, size=12.5)
add_text(s, 0.6, 5.35, 12.1, 0.4, "作品定位：给独居老人配一个“看得见、记得住、靠得住”的随身智能体。", size=16, bold=True, color=AMBER)
add_text(s, 0.6, 5.9, 12.1, 0.4, "为什么是眼镜：找物、确认服药、判断跌倒都需要第一视角——手机和手表做不到；语音符合老人习惯；免去主动操作。", size=13, color=MUTED)
add_footer(s, 3)
set_notes(s, "重点说明：作品不是泛问答，而是感知驱动 + 主动触发 + 记忆沉淀的原生智能体。")

# ================= 04 目标用户与痛点 =================
s = new_slide()
add_header(s, 4, "2 目标用户与场景痛点", "60–80 岁独居 / 白天独居、轻度记忆衰退、生活基本自理的老人")
cards = [
    ("遗忘", "找不到老花镜、钥匙\n忘记吃药、重复吃药\n反复翻找、着急无助", RED),
    ("安全", "跌倒无人知晓\n出门忘带伞、忘关火\n错过黄金救援时间", AMBER),
    ("孤独", "子女不在身边\n没人说话、情绪低落\n子女打电话追问反而嫌烦", PER),
    ("子女焦虑", "“吃药了吗？”式电话\n问了嫌烦、不问担心\n无法获得日常真实状态", GREEN),
]
cw = (12.13 - 0.3) / 4
for i, (t, d, c) in enumerate(cards):
    box = add_rect(s, 0.6 + i * (cw + 0.1), 1.7, cw, 3.1, fill=PANEL, line=c)
    shape_text(box, f"{t}\n\n{d}", size=13.5)
rows = [
    ("一级用户", "独居老人：轻度健忘、生活自理、不愿用手机"),
    ("二级用户", "子女（远程安心）与社区护工（批量服务）"),
    ("典型画像", "王奶奶，68 岁，独居，轻度健忘：常找不到老花镜、偶尔忘吃降压药"),
]
for i, (t, d) in enumerate(rows):
    add_text(s, 0.6, 5.1 + i * 0.5, 1.7, 0.4, t, size=13.5, bold=True, color=ACCENT)
    add_text(s, 2.5, 5.1 + i * 0.5, 10.1, 0.4, d, size=13)
add_footer(s, 4)
set_notes(s, "痛点围绕记忆衰退老人的三大困境，用户分层明确。")

# ================= 05 需求分析 =================
s = new_slide()
add_header(s, 5, "3 需求分析", "功能需求 + 非功能需求")
add_text(s, 0.6, 1.5, 12.1, 0.4, "功能需求（P0 = 可演示闭环，P1 = 安全守护与陪伴）", size=15, bold=True)
reqs = [
    ("P0", "记忆找物", "语音提问 → 记忆检索 → 回答 → 复核更新"),
    ("P0", "用药提醒+确认", "时间触发 → 主动提醒 → 视觉确认 → 健康日志"),
    ("P0", "出门提醒", "场景触发 → 天气+室内状态联动提醒"),
    ("P1", "跌倒检测求助", "IMU+视觉双通道 → 二次确认 → 无回应联系家人"),
    ("P1", "主动陪伴 / 健康日报", "低打扰陪伴，每晚给家人一条摘要"),
    ("P1", "一键紧急求助", "镜腿长按 / 语音 SOS → 联系家人 + 位置"),
]
for i, (pri, t, d) in enumerate(reqs):
    r, c = divmod(i, 3)
    box = add_rect(s, 0.6 + c * 4.12, 2.0 + r * 1.15, 3.92, 1.0, fill=PANEL, line=ACCENT if pri == "P0" else BORDER, radius=0.1)
    shape_text(box, f"{pri} · {t}\n{d}", size=12)
add_text(s, 0.6, 4.5, 12.1, 0.4, "非功能需求", size=15, bold=True)
nfr = [
    ("低打扰", "只在值得开口时主动说话，有优先级评估"),
    ("可解释", "每次决策可回溯：感知→触发→工具→反馈"),
    ("隐私", "本地优先、最小必要、可查看可删除"),
    ("可用性", "无 Key 时给出明确引导；断网有兜底"),
]
for i, (t, d) in enumerate(nfr):
    box = add_rect(s, 0.6 + i * 3.06, 5.0, 2.86, 1.3, fill=PANEL2, line=BORDER, radius=0.1)
    shape_text(box, f"{t}\n{d}", size=12)
add_footer(s, 5)
set_notes(s, "需求分 P0/P1：本次比赛完成 P0 三个闭环，P1 展示完整产品边界。")

# ================= 06 功能设计（一） =================
s = new_slide()
add_header(s, 6, "4 功能设计（一）核心功能", "F1–F3：三个可演示闭环")
funcs = [
    ("F1 记忆找物", "被动响应 · 语音提问", "检索第一视角记忆中的物品位置，随问随答；找到后复核更新置信度。记忆来源：日常自然“看到”的高价值物品，不做持续录像。", "P0"),
    ("F2 用药提醒+确认", "主动服务 · 时间触发", "到点主动提醒；视觉识别药盒确认；确认后写入健康日志；10/30 分钟未确认逐级升级通知家人。不做剂量建议，不替代医嘱。", "P0"),
    ("F3 出门提醒", "主动服务 · 场景触发", "识别走向门口时，联动天气与室内状态（灶台火）提醒；确认关火后更新状态，避免重复唠叨。只提醒，不控制门锁/燃气。", "P0"),
]
for i, (t, trig, d, pri) in enumerate(funcs):
    box = add_rect(s, 0.6, 1.7 + i * 1.7, 12.13, 1.5, fill=PANEL2, line=ACCENT, radius=0.08)
    add_text(s, 0.9, 1.88 + i * 1.7, 3.2, 0.4, t, size=16, bold=True, color=ACCENT)
    add_text(s, 0.9, 2.3 + i * 1.7, 3.2, 0.35, trig, size=12, color=AMBER)
    add_text(s, 4.3, 1.9 + i * 1.7, 8.1, 1.2, d, size=12.5)
    add_text(s, 11.9, 1.8 + i * 1.7, 0.7, 0.4, pri, size=13, bold=True, color=GREEN)
add_footer(s, 6)
set_notes(s, "三个 P0 闭环是本次作品的核心，均给出边界（不做诊断、不做控制）。")

# ================= 07 功能设计（二） =================
s = new_slide()
add_header(s, 7, "4 功能设计（二）扩展功能与反功能", "F4–F7 安全守护与陪伴")
ext = [
    ("F4 跌倒检测求助", "异常触发", "IMU 姿态突变 + 视觉人形倒地双通道确认；语音二次确认；无回应 15 秒自动联系家人，附定位+快照。"),
    ("F5 主动陪伴/播报", "时间/情绪触发", "长时间无对话或情绪低落时主动关怀；天气/新闻播报、回忆话题。涉及健康建议时引导就医。"),
    ("F6 家人健康日报", "每日汇总", "每晚 8 点推送给子女：服药情况、出门时段、异常事件；一切正常只发“今天都挺好”。"),
    ("F7 一键紧急求助", "被动响应", "镜腿长按 3 秒或语音 SOS → 呼叫紧急联系人 + 发送位置。"),
]
for i, (t, trig, d) in enumerate(ext):
    r, c = divmod(i, 2)
    box = add_rect(s, 0.6 + c * 6.16, 1.7 + r * 1.85, 5.96, 1.6, fill=PANEL, line=BORDER, radius=0.08)
    add_text(s, 0.85 + c * 6.16, 1.85 + r * 1.85, 4.5, 0.4, t, size=15, bold=True)
    add_text(s, 0.85 + c * 6.16, 2.25 + r * 1.85, 2.6, 0.35, trig, size=11.5, color=AMBER)
    add_text(s, 0.85 + c * 6.16, 2.62 + r * 1.85, 5.5, 0.7, d, size=11.5, color=MUTED)
box = add_rect(s, 0.6, 5.5, 12.13, 1.1, fill=PANEL, line=RED, radius=0.08)
shape_text(box, "明确不做（反功能）\n不做医疗诊断 · 不持续云端录像 · 不识别陌生人面孔 · 不做广告营销打扰", size=13)
add_footer(s, 7)
set_notes(s, "F4–F7 展示完整产品边界；反功能设计是合规与产品成熟度的体现。")

# ================= 08-11 核心闭环 =================
flow_slide(
    8, "5 核心任务闭环", "闭环一 · 记忆找物", "第一视角记忆：看过即记住，随问随答，找到后复核更新",
    [("① 语音提问", "“我的老花镜呢？”"), ("② 检索记忆", "物品-位置-置信度"), ("③ 语音回答", "“在客厅茶几上”"), ("④ 找到确认", "“找到了”"), ("⑤ 复核更新", "记忆置信度刷新")],
    [("记忆未命中", "现场视觉扫描高频区域，找到后回写", AMBER), ("没有把握", "明确告知，不编造位置", RED)],
    "演示实录：09:12 提问 → 记忆库检索命中（置信度 0.9）→ 语音回答 → 09:13 复核更新。",
)
flow_slide(
    9, "5 核心任务闭环", "闭环二 · 用药提醒（主动 Agent）", "时间触发 → 主动开口 → 视觉确认 → 健康日志 → 漏服升级",
    [("① 08:00 触发", "用药计划到点"), ("② 主动开口", "“该吃降压药了”"), ("③ 视觉确认", "识别到药盒"), ("④ 确认服药", "“吃完了”"), ("⑤ 写入日志", "健康日志+日报")],
    [("10 分钟未确认", "震动 + 再次提醒", AMBER), ("30 分钟未确认", "通知子女“可能漏服”", RED)],
    "演示实录：08:00 主动开口（时间触发 HIGH）→ 08:02 视觉确认 → 08:04 写入服药记录。",
)
flow_slide(
    10, "5 核心任务闭环", "闭环三 · 出门提醒", "场景触发 → 天气 API + 室内状态联动 → 确认后状态更新",
    [("① 走向门口", "识别出门意图"), ("② 状态查询", "天气 + 灶台"), ("③ 联动提醒", "“带伞，火没关”"), ("④ 返回关火", "“放心出门”"), ("⑤ 状态更新", "记忆刷新")],
    [("再次出门", "只提醒天气，不重复唠叨", GREEN), ("30 分钟未归", "家人低打扰通知", AMBER)],
    "演示实录：08:50 场景触发 → 08:52 确认关火更新状态 → 08:54 再次出门只提醒天气。",
)
flow_slide(
    11, "5 核心任务闭环", "闭环四 · 跌倒检测与求助（安全守护）", "IMU + 视觉双通道确认，无回应自动联系家人",
    [("① 双通道检测", "姿态突变+人形倒地"), ("② 语音二次确认", "“您摔着了吗？”"), ("③ 有回应", "通知家人一条"), ("④ 无回应 15s", "自动呼叫联系人"), ("⑤ 事件沉淀", "健康日志+误报样本")],
    [("责任边界", "提醒+通知工具，非医疗设备", MUTED), ("家属联动", "附定位 + 第一视角快照", PER)],
    "安全守护是产品完整性加分项：双通道降误报，二次确认避免打扰。",
)

# ================= 12 Agent 设计 =================
s = new_slide()
add_header(s, 12, "6 Agent 设计 · 原生智能体", "主动触发引擎：不是泛问答，而是自主判断“是否值得开口”")
triggers = [
    ("⏰ 时间触发", "08:00 用药计划", "“该吃降压药了”", "HIGH"),
    ("🚪 场景触发", "识别到走向门口", "“下雨带伞，火没关”", "HIGH"),
    ("🆘 异常触发", "跌倒/久未活动", "“您摔着了吗？”", "HIGH"),
    ("💡 记忆触发", "重复问同一物品", "“您上次问过…”", "MEDIUM"),
    ("💬 情绪触发", "语气低沉/久未说话", "“出去走走吗？”", "LOW"),
]
cw = (12.13 - 0.4) / 5
for i, (t, cond, out, pri) in enumerate(triggers):
    box = add_rect(s, 0.6 + i * (cw + 0.1), 1.6, cw, 2.7, fill=PANEL2, line=BORDER)
    shape_text(box, f"{t}\n\n条件：{cond}\n\n输出：{out}\n\n优先级 {pri}", size=11.5)
box = add_rect(s, 0.6, 4.55, 12.13, 1.15, fill=PANEL, line=ACCENT, radius=0.08)
shape_text(box, "决策守卫：是否开口 = f(事件价值, 紧急程度, 打扰成本, 当前上下文)\n深夜不闲聊 · 跌倒立即触发 · 对话中降级为静默记录", size=13.5)
box = add_rect(s, 0.6, 5.9, 12.13, 0.85, fill=PANEL, line=GREEN, radius=0.08)
shape_text(box, "可靠性兜底：用药时间 / 跌倒等关键触发由系统规则保证（规则+LLM 混合），演示不依赖模型偶然性", size=13)
add_footer(s, 12)
set_notes(s, "原生智能体的核心证据：主动性 + 优先级 + 可回溯 + 规则兜底。")

# ================= 13 Agent 架构 =================
s = new_slide()
add_header(s, 13, "6 Agent 设计 · 分层架构", "感知 → 情境 → 决策 → 工具 → 反馈 → 记忆")
layers = [
    ("感知层", "摄像头帧流（第一视角）· 语音输入 · 场景描述", ACCENT),
    ("情境引擎", "场景 / 人物状态 / 可见物体 / 时间 / 天气上下文", PER),
    ("Agent 核心", "意图理解 · 主动触发引擎 · 决策 · 规则兜底", AMBER),
    ("工具层", "记忆检索/更新 · 用药计划 · 天气 · 家人通知", GREEN),
    ("反馈层", "语音反馈（一句话）· 决策轨迹可视化", RED),
    ("记忆沉淀", "物品位置 · 健康日志 · 事件 → 本地 JSON 持久化", MUTED),
]
y = 1.42
bh = 0.72
for i, (t, d, c) in enumerate(layers):
    box = add_rect(s, 0.6, y, 12.13, bh, fill=PANEL2, line=c, radius=0.08)
    add_text(s, 0.85, y + 0.16, 2.1, 0.4, t, size=15, bold=True, color=c)
    add_text(s, 3.05, y + 0.18, 9.5, 0.42, d, size=12.5)
    if i < len(layers) - 1:
        add_arrow(s, 6.35, y + bh + 0.015, 0.45, 0.17, color=BORDER, direction="down")
    y += bh + 0.21
add_text(s, 0.6, 6.45, 12.1, 0.4, "决策全程可回溯 → 隐私审计 + 演示可视化（决策轨迹是“原生智能体”的直接证据）", size=13, bold=True, color=AMBER)
add_footer(s, 13)
set_notes(s, "真实版架构：规则兜底保证可靠性，LLM 负责理解与表达，工具真实执行。")

# ================= 14 记忆系统 =================
s = new_slide()
add_header(s, 14, "6 Agent 设计 · 记忆系统", "记忆即能力：找物、用药、出门提醒都依赖持续记忆")
mems = [
    ("短期记忆", "当前对话、最近交互", "上下文连续性"),
    ("长期记忆", "物品-位置（时间戳+置信度）、用药记录、健康事件", "随问随答"),
    ("遗忘策略", "置信度降权 · 周期性清理 · 支持“忘记这个”", "隐私友好"),
    ("家人可见", "云端仅同步摘要（服药/出门/异常）", "低打扰"),
]
for i, (t, d, tag) in enumerate(mems):
    box = add_rect(s, 0.6 + (i % 2) * 6.16, 1.6 + (i // 2) * 1.8, 5.96, 1.55, fill=PANEL, line=ACCENT if i == 1 else BORDER)
    shape_text(box, f"{t}\n{d}", size=13)
add_text(s, 0.6, 5.5, 12.1, 0.4, "隐私设计：本地优先 · 最小必要 · 高价值物品才记忆 · 可查看与删除（前端可点击删除）", size=15, bold=True, color=AMBER)
add_text(s, 0.6, 6.0, 12.1, 0.4, "真实版实现：记忆持久化到 app/data/memory.json，刷新不丢", size=13, color=MUTED)
add_footer(s, 14)
set_notes(s, "记忆是本作品的核心资产，同时把隐私设计成产品特性。")

# ================= 15 技术实现 =================
s = new_slide()
add_header(s, 15, "7 技术实现 · 真实可用版（app/）", "零第三方依赖：Python 标准库 HTTP 服务 + 浏览器前端")
techs = [
    ("服务端", "Python 标准库 http.server，零依赖；接口：/api/run、/api/config、/api/state、/api/test 等"),
    ("LLM 接入", "OpenAI 兼容接口：OpenAI / DeepSeek / Moonshot / 通义千问 / 本地 Ollama；支持视觉 + JSON 输出"),
    ("视觉感知", "摄像头拍照 / 上传照片 → 视觉模型识别场景、人物状态、可见物体；不支持视觉时用手动场景描述"),
    ("语音", "浏览器 Web Speech API：语音识别输入 + 语音朗读反馈（无需额外 ASR/TTS Key）"),
    ("工具层", "记忆检索/更新、用药计划、天气（可选 Key）、家人通知（可选 Webhook）真实执行"),
    ("记忆持久化", "物品位置、健康日志保存到 app/data/memory.json，可查看、可删除"),
]
for i, (t, d) in enumerate(techs):
    box = add_rect(s, 0.6, 1.5 + i * 0.88, 12.13, 0.72, fill=PANEL, line=BORDER, radius=0.1)
    add_text(s, 0.85, 1.66 + i * 0.88, 1.9, 0.4, t, size=14, bold=True, color=ACCENT)
    add_text(s, 2.9, 1.66 + i * 0.88, 9.6, 0.45, d, size=12.5)
add_footer(s, 15)
set_notes(s, "真实版强调可落地：用户只提供 AI Key，应用即可运行。")

# ================= 16 系统使用说明 =================
s = new_slide()
add_header(s, 16, "7 技术实现 · 系统使用说明", "三步启动，浏览器操作")
steps = [
    ("① 启动服务", "cd app\npy server.py", "本机 http://127.0.0.1:8000"),
    ("② 配置 Key", "「配置」页填写 API Key / 模型\n保存 → 测试连接", "Key 仅存本机 config.json"),
    ("③ 使用", "摄像头/照片/手动场景\n语音或文字输入 → 运行一步", "三个闭环均可操作"),
]
for i, (t, d, note) in enumerate(steps):
    box = add_rect(s, 0.6 + i * 4.12, 1.7, 3.92, 2.2, fill=PANEL2, line=ACCENT, radius=0.08)
    shape_text(box, f"{t}\n\n{d}", size=13)
    add_text(s, 0.6 + i * 4.12, 4.0, 3.92, 0.4, note, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2:
        add_arrow(s, 0.6 + (i + 1) * 4.12 - 0.28, 2.65, 0.26, 0.3, color=MUTED, direction="right")
add_text(s, 0.6, 4.75, 12.1, 0.4, "前端操作台能力：", size=15, bold=True)
feats = [
    "📷 摄像头拍照 / 上传图片（第一视角）",
    "🎤 浏览器语音输入 + 语音朗读",
    "⚡ 快捷场景与快捷指令",
    "⏰ 模拟 08:00 用药时间（演示主动提醒）",
    "📋 决策轨迹实时展示",
    "🧠 记忆面板实时更新（可删除）",
]
for i, f in enumerate(feats):
    r, c = divmod(i, 2)
    add_text(s, 0.6 + c * 6.16, 5.3 + r * 0.5, 5.9, 0.4, f, size=12.5)
add_text(s, 0.6, 6.35, 12.13, 0.5,
         "说明：网页端为初赛演示形态（便于评委打开浏览器复现）；正式产品形态为「AI 眼镜 + 老人端 App + 子女端 App」，见下两页。",
         size=12.5, color=AMBER)
add_footer(s, 16)
set_notes(s, "使用门槛低：不需要硬件、不需要部署云端，本机一条命令启动。网页端是初赛演示形态。")

# ================= 17 产品形态（一）老人端 =================
s = new_slide()
add_header(s, 17, "7 技术实现 · 产品形态：老人端 App 构想",
           "网页端为初赛演示形态；正式产品为「AI 眼镜 + 老人端 App + 子女端 App」")
forms = [
    ("初赛：Web 演示", "app/ + web/\n浏览器打开即用", ACCENT),
    ("正式：AI 眼镜", "语音 + 第一视角\n随身主入口", GREEN),
    ("老人端 App", "适老化大字界面\n设置 / SOS / 记忆", PER),
    ("子女端 App", "远程日报 / 通知\n留言 / 看护", AMBER),
]
fw = (12.13 - 0.9) / 4
for i, (t, d, c) in enumerate(forms):
    box = add_rect(s, 0.6 + i * (fw + 0.3), 1.42, fw, 1.05, fill=PANEL2, line=c, radius=0.08)
    shape_text(box, f"{t}\n{d}", size=12)
    if i < 3:
        add_arrow(s, 0.6 + (i + 1) * (fw + 0.3) - 0.27, 1.8, 0.24, 0.28, color=MUTED, direction="right")
add_text(s, 0.6, 2.72, 12.13, 0.35,
         "老人端设计原则：大字大按钮（≥20pt / ≥56pt）· 语音优先 · 层级≤2 · 防误触二次确认 · 支持方言",
         size=13, bold=True, color=AMBER)
app_funcs = [
    ("首页 · 今日状态", "大字时间/天气/服药状态\n一键“暖眸播报”"),
    ("语音助手", "全双工对话\n常用问题快捷卡"),
    ("找物记录", "物品位置清单\n“我的东西都在哪”"),
    ("用药管理", "计划卡片 + 已服/待服\n漏服提醒与记录"),
    ("安全守护", "跌倒记录 / 当前位置\nSOS 说明与联系人"),
    ("亲情留言", "子女语音留言播放\n家庭相册回忆"),
]
gw = (12.13 - 0.3) / 3
for i, (t, d) in enumerate(app_funcs):
    r, c = divmod(i, 3)
    box = add_rect(s, 0.6 + c * (gw + 0.15), 3.15 + r * 1.1, gw, 0.95, fill=PANEL, line=BORDER, radius=0.1)
    shape_text(box, f"{t}\n{d}", size=11.5)
box = add_rect(s, 0.6, 5.5, 12.13, 1.35, fill=PANEL2, line=ACCENT, radius=0.06)
shape_text(box,
           "首页界面构想（线框）：顶部 大字时间 + 天气 ┃ 中部 用药大卡片（今日已服 ✓ / 待服）┃ "
           "底部 三大金刚按钮：语音助手 · 找东西 · SOS（红色大字）",
           size=12.5)
add_footer(s, 17)
set_notes(s, "老人端是眼镜的伴侣屏：大屏展示、设置、SOS 与记忆查看；交互以语音优先，严格适老化。")

# ================= 18 产品形态（二）子女端 =================
s = new_slide()
add_header(s, 18, "7 技术实现 · 产品形态：子女端 App 与双端联动", "子女端：远程守护与低打扰沟通")
family_funcs = [
    ("今日概况日报", "“今天都挺好”\n异常才展开详情"),
    ("实时状态", "位置 / 眼镜电量\n是否佩戴 / 活动"),
    ("用药看板", "今日与本周记录\n漏服红色提醒"),
    ("异常通知", "跌倒 / 久未活动\n深夜未归 + 快照"),
    ("沟通", "语音留言给老人\n视频通话入口"),
    ("隐私与偏好", "可见范围 / 通知级别\n低打扰设置"),
]
gw = (12.13 - 0.3) / 3
for i, (t, d) in enumerate(family_funcs):
    r, c = divmod(i, 3)
    box = add_rect(s, 0.6 + c * (gw + 0.15), 1.5 + r * 1.1, gw, 0.95, fill=PANEL, line=BORDER, radius=0.1)
    shape_text(box, f"{t}\n{d}", size=11.5)
add_text(s, 0.6, 3.95, 12.13, 0.35, "双端与眼镜联动（数据流）", size=14, bold=True)
flow = [
    ("AI 眼镜", "感知 / 语音 / 反馈", ACCENT),
    ("老人端 App", "大字显示 / SOS / 设置", PER),
    ("云端", "记忆 + 健康日志摘要", MUTED),
    ("子女端 App", "日报 / 通知 / 留言", AMBER),
]
fw = (12.13 - 0.6) / 4
for i, (t, d, c) in enumerate(flow):
    box = add_rect(s, 0.6 + i * (fw + 0.2), 4.4, fw, 0.8, fill=PANEL2, line=c, radius=0.08)
    shape_text(box, f"{t}\n{d}", size=11.5)
    if i < 3:
        add_arrow(s, 0.6 + (i + 1) * (fw + 0.2) - 0.2, 4.66, 0.2, 0.26, color=MUTED, direction="right")
add_text(s, 0.6, 5.5, 6.4, 0.6,
         "首页界面构想：老人状态卡（头像 + 今日摘要）\n底部 Tab：首页 / 用药 / 安全 / 我的", size=12.5)
add_text(s, 7.1, 5.5, 5.7, 0.6,
         "通知分级：正常一句“今天都挺好”\n提醒（漏服 / 久未活动）· 紧急（跌倒红色）", size=12.5)
add_text(s, 0.6, 6.4, 12.13, 0.4,
         "隐私边界：子女仅可见授权摘要（服药 / 出门 / 异常）；完整记忆留在本机，可删除", size=12.5, color=MUTED)
add_footer(s, 18)
set_notes(s, "子女端价值：低打扰守护——正常不打扰，异常才提醒；所有数据按授权最小可见。")

# ================= 19 数据与合规 =================
s = new_slide()
add_header(s, 19, "8 数据来源与合规边界", "数据授权 · 隐私保护 · 风险提示 · 行业边界")
comps = [
    ("数据来源", "用户与家属授权采集\n开源/自采模拟数据\n演示用剧本数据", PER),
    ("隐私保护", "本地优先 · 最小必要\n人脸模糊 · 加密存储\n可查看、可删除、可撤回", GREEN),
    ("合规边界", "个人信息保护法\n采集指示灯 + 知情同意\n适老化认证", AMBER),
    ("风险声明", "不做医疗诊断（提醒工具）\n跌倒求助为通知辅助\n误报/依赖/误触发有预案", RED),
]
cw4 = (12.13 - 0.3) / 4
for i, (t, d, c) in enumerate(comps):
    box = add_rect(s, 0.6 + i * (cw4 + 0.1), 1.7, cw4, 3.4, fill=PANEL, line=c)
    shape_text(box, f"{t}\n\n{d}", size=12.5)
add_text(s, 0.6, 5.4, 12.1, 0.4, "核心原则：服务的是“提醒与陪伴”，把医疗诊断、急救责任明确挡在边界之外。", size=15, bold=True, color=AMBER)
add_text(s, 0.6, 5.95, 12.1, 0.4, "真实版实现：API Key 仅保存在本机 app/config.json（已 gitignore，不会进入代码仓库）", size=13, color=MUTED)
add_footer(s, 19)
set_notes(s, "合规是评审必问项：数据从哪来、存哪里、谁能删、出问题谁负责。")

# ================= 20 Demo 与验证 =================
s = new_slide()
add_header(s, 20, "9 Demo 与验证", "真实运行演示 + 自动化测试 + 无 Key 兜底")
box = add_rect(s, 0.6, 1.6, 5.9, 3.4, fill=PANEL, line=ACCENT, radius=0.06)
shape_text(box, "真实运行演示（app/）\n\n① 摄像头/照片 → 视觉识别\n② 三个闭环逐步运行\n③ 决策轨迹实时展示\n④ 记忆面板实时变化\n\n浏览器即可复现", size=13)
box = add_rect(s, 6.75, 1.6, 5.98, 3.4, fill=RGBColor(0x0B, 0x0E, 0x1C), line=BORDER, radius=0.06)
shape_text(box, "[感知] 08:00 场景=卧室 · 语音=无\n[触发] ⏰ 时间触发 · 用药计划（HIGH）\n[决策] 动作=remind_medication\n[工具] 用药计划 → 降压药（1粒/日）\n[反馈] 🔊 王奶奶，该吃降压药了。\n\n= Agent 决策轨迹（真实执行）", size=12, align=PP_ALIGN.LEFT)
checks = [
    ("端到端测试", "Mock LLM 模拟服务商：12 项断言全部通过（视觉→决策→工具→记忆）"),
    ("HTTP 冒烟测试", "前端页面、JS 类型、健康检查、无 Key 兜底提示均正常"),
    ("无 Key 兜底", "未配置 Key 时给出明确引导，不崩溃"),
    ("概念演示版", "web/（浏览器版）与 demo/（Python 版）无需 Key，可作视频兜底材料"),
]
for i, (t, d) in enumerate(checks):
    r, c = divmod(i, 2)
    box = add_rect(s, 0.6 + c * 6.16, 5.3 + r * 0.72, 5.96, 0.6, fill=PANEL2, line=BORDER, radius=0.1)
    add_text(s, 0.85 + c * 6.16, 5.42 + r * 0.72, 5.5, 0.45, f"{t}：{d}", size=11.5)
add_footer(s, 20)
set_notes(s, "验证材料：可运行代码 + 自动化测试 + 演示视频 + 在线仓库，评审可自行复现。")

# ================= 21 迭代计划 =================
s = new_slide()
add_header(s, 21, "10 迭代计划", "从初赛可运行版本到长期陪伴生态")
vers = [
    ("V1 · 初赛交付", "真实 Agent 版（app/）\n找物/用药/出门三闭环\nWeb 概念演示 + 材料", ACCENT),
    ("V2 · 安全守护", "跌倒检测求助\n家人健康日报\n心率/血氧体征", GREEN),
    ("V3 · 规模服务", "方言扩展\n社区护工多对多\n端侧模型离线化", AMBER),
    ("长期 · 生态", "可复用模板\n智能体商店\n开放工具组件", RED),
]
cw5 = (12.13 - 0.9) / 4
for i, (t, d, c) in enumerate(vers):
    box = add_rect(s, 0.6 + i * (cw5 + 0.3), 1.7, cw5, 2.7, fill=PANEL2, line=c)
    shape_text(box, f"{t}\n\n{d}", size=13.5)
    if i < 3:
        add_arrow(s, 0.6 + (i + 1) * (cw5 + 0.3) - 0.26, 2.85, 0.24, 0.3, color=MUTED, direction="right")
add_text(s, 0.6, 4.95, 12.1, 0.5, "设计原则：先做深一个场景，再做宽多个场景；每一步都保留可演示、可验证的任务闭环。", size=15, bold=True, color=AMBER)
add_footer(s, 21)
set_notes(s, "V1 是本次初赛交付，后续版本保持闭环可验证原则。")

# ================= 22 交付物与结尾 =================
s = new_slide()
add_header(s, 22, "10 迭代计划 · 交付物清单", "参赛提交材料")
deliverables = [
    ("真实可用版", "app/：零依赖 Agent 应用，接入 AI Key 即用", GREEN),
    ("概念演示版", "web/ 浏览器版 + demo/ Python 版（无需 Key）", PER),
    ("设计文档", "docs/：产品简介、产品设计方案、功能闭环演示", ACCENT),
    ("项目说明 PPT", "本文件（20 页，文档向）", AMBER),
    ("打包与仓库", "zip 提交包 + GitHub 仓库（AIglasses）", RED),
    ("验证材料", "端到端测试 + 功能闭环演示（场景化）", MUTED),
]
for i, (t, d, c) in enumerate(deliverables):
    r, cc = divmod(i, 2)
    box = add_rect(s, 0.6 + cc * 6.16, 1.7 + r * 0.95, 5.96, 0.8, fill=PANEL, line=c, radius=0.1)
    add_text(s, 0.85 + cc * 6.16, 1.86 + r * 0.95, 1.9, 0.4, t, size=14, bold=True, color=c)
    add_text(s, 2.85 + cc * 6.16, 1.86 + r * 0.95, 3.6, 0.5, d, size=11.5)
add_text(s, 0.9, 5.0, 11.5, 0.8, "暖眸——看得见、记得住、靠得住", size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 5.9, 11.5, 0.5, "让每一位独居老人，都有一位 24 小时在身边的智能守护者", size=15, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 6.55, 11.5, 0.4, "谢谢 · 请各位评委指正", size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 22)
set_notes(s, "结尾汇总交付物，强调可复现、可验证。")

# ---------------- 保存 ----------------
import argparse
import os

_parser = argparse.ArgumentParser(description="生成暖眸项目说明 PPT")
_parser.add_argument(
    "--out",
    default=os.path.join(os.path.dirname(os.path.abspath(__file__)), "暖眸-居家养老AI眼镜智能体-项目说明.pptx"),
    help="输出 pptx 路径",
)
_args = _parser.parse_args()

prs.save(_args.out)
print("已生成：", _args.out)
print("页数：", len(prs.slides))
